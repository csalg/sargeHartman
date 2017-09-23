# Allow user to input directory
# Commit to git

# To dos
# Use decorator
# Improve theme




from pptx import Presentation
import os
from random import randint
from PIL import Image
from resizeimage import resizeimage

try:
	os.mkdir('img_res')
except Exception as excpt:
	print(excpt)


def resizeImg(img):
	with open(os.path.join('img', img), 'r+b') as f:
	    with Image.open(f) as image:
	        cover = resizeimage.resize_cover(image, [1024, 800])
	        cover.save(os.path.join('img_res',img), image.format)
# Variables

template = Presentation('template.pptx') # Create a new presentation object
prs = Presentation()
wordsFiles = os.listdir('img')
words = [word[:word.find('.')] for word in wordsFiles]
if '' in words:
	wordsFiles.pop(words.index(''))
	words.pop(words.index(''))

for image in wordsFiles:
	resizeImg(image)

# Functions

def titleSlide(text):
	global prs
	emptyprs = Presentation()
	slideTitleLayout = emptyprs.slide_layouts[0] # This grabs the layout from the first slide
	slide = prs.slides.add_slide(slideTitleLayout) # Creates a new slide based on the layout we grabbed.
	slide.shapes.title.text = text

def emuer(inches):
	return int(float(inches)*914400.0)

def twoColumns(word1,word2):
	global prs
	slide = prs.slides.add_slide(template.slides[1].slide_layout)
	slide.shapes[1].text = word1
	slide.shapes[1].left = emuer(1.88)
	slide.shapes[1].top = emuer(1.1)
	slide.shapes[2].text = word2
	slide.shapes[2].left = emuer(5.56)
	slide.shapes[2].top = emuer(1.1)

# Use a decorator
def pogoDrill(word, other_words):
	side = randint(0,1)
	if side == 0:
		return twoColumns(other_words[randint(0,len(other_words)-1)], word)
	else:
		return twoColumns(word, other_words[randint(0,len(other_words)-1)])

def phonicsDriller(word):
	print(words)
	print(word)
	global prs
	titleSlide(word)
	slide = prs.slides.add_slide(prs.slide_layouts[6])
	print(f'Creating image slide for {word}')
	slide.shapes.add_picture(os.path.join('img_res', wordsFiles[words.index(word)]), 0, 0)
	titleSlide(word)

#  <------------------------ Control flow ------------------------->

# Letters and phonics

titleSlide('Scaffolded spelling drills')
for word in words:
	phonicsDriller(word)

# Pogo drill
titleSlide('Pogo drill')
for word in words:
	titleSlide(word)
	other_words= words[:]
	other_words.pop(words.index(word))
	print(word)
	print(other_words)
	for _ in range(1,10):
		pogoDrill(word, other_words)

prs.save('a.pptx')
